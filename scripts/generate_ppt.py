"""
Generate group meeting PPT v2 — with all 10 improvements.
"""
import json, os, glob as _glob
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
ASSETS = os.path.join(ROOT, 'assets', 'figures')
OUTPUT = os.path.join(ROOT, 'docs', '组会汇报.pptx')
os.makedirs(ASSETS, exist_ok=True)

# ─── Chinese font for matplotlib ───
_font_candidates = ['C:/Windows/Fonts/msyh.ttc','C:/Windows/Fonts/simsun.ttc','C:/Windows/Fonts/simhei.ttf']
_font_candidates += _glob.glob('C:/Windows/Fonts/*.ttf') + _glob.glob('C:/Windows/Fonts/*.ttc')
_zh_name = 'sans-serif'
for _fp in _font_candidates:
    if os.path.exists(_fp):
        from matplotlib.font_manager import FontProperties as _FP
        _zh_name = _FP(fname=_fp).get_name()
        break
plt.rcParams.update({
    'font.size':12,'font.family':'sans-serif','font.sans-serif':[_zh_name,'DejaVu Sans'],
    'axes.titlesize':14,'axes.labelsize':12,'axes.unicode_minus':False,
    'figure.facecolor':'white','axes.facecolor':'white','axes.edgecolor':'#333333',
    'axes.grid':True,'grid.alpha':0.3,'axes.spines.top':False,'axes.spines.right':False,
})

# ─── Colors ───
DARK_BG  = RGBColor(0x1A,0x1A,0x2E)
PRIMARY  = RGBColor(0x00,0x9E,0xFF)
ORANGE   = RGBColor(0xFF,0x6B,0x35)
GREEN    = RGBColor(0x00,0xCC,0x88)
RED      = RGBColor(0xFF,0x41,0x41)
WHITE    = RGBColor(0xFF,0xFF,0xFF)
GRAY     = RGBColor(0xCC,0xCC,0xCC)
LIGHT_GRAY=RGBColor(0xF0,0xF0,0xF0)
BLACK    = RGBColor(0x20,0x20,0x20)
YELLOW   = RGBColor(0xFF,0xC1,0x07)
PURPLE   = RGBColor(0x9C,0x27,0xB0)

def _to_rgb(c):
    if isinstance(c,RGBColor): return c
    if isinstance(c,str) and c.startswith('#'):
        c=c.lstrip('#')
        return RGBColor(int(c[0:2],16),int(c[2:4],16),int(c[4:6],16))
    return c

# ═══════════════════════ Data ═══════════════════════
with open(os.path.join(ROOT,'experiments','phase3_ablation.json')) as f:
    phase3 = json.load(f)
with open(os.path.join(ROOT,'experiments','cross_language_results.json')) as f:
    xl_data = json.load(f)

experiments = ['A1\n(基线)','A2\n(统计先验)','A2b\n(随机初始化)','B3\n(Adapter+Prosody)','A3\n(Full Fine-tune)']
test_means = []
for key in ['A1_baseline','A2_stat_prior','A2b_rand_init','B3_prosody','A3_full_finetune']:
    vals = [phase3[key][s]['test'] for s in ['seed_42','seed_123','seed_456']]
    test_means.append(sum(vals)/len(vals))
fd_exp = ['C1(无增强)','C3(儿童约束)','C2(成人参数)','C4(极端参数)']
fd_fd = [0,8.71,9.87,11.99]
fd_acc = [86.42,66.85,59.51,43.25]
old_models = ['DrseCNN','CNN','BiLSTM','Transformer']
old_before = [85.99,72.14,72.26,69.48]
old_after  = [35,30,30,28]

# ═══════════════════════ Chart functions ═══════════════════════
def save_chart(fig, name):
    path = os.path.join(ASSETS, f'chart_{name}.png')
    fig.savefig(path, dpi=150, bbox_inches='tight')
    plt.close(fig)
    return path

def chart_fd_vs_acc():
    fig,ax=plt.subplots(figsize=(8,4.5))
    ax.plot(fd_fd,fd_acc,'o-',color='#FF4136',linewidth=2,markersize=10,markerfacecolor='white',markeredgewidth=2)
    for x,y,label in zip(fd_fd,fd_acc,fd_exp):
        offset=2 if x==0 else -3
        ax.annotate(f'{label} ({y}%)',(x,y),textcoords="offset points",xytext=(10,offset),fontsize=10,color='#333333')
    ax.set_xlabel('Frechet Distance (FD)'); ax.set_ylabel('Validation Accuracy (%)')
    ax.set_title('分布偏移量 vs 分类性能：严格负相关')
    fig.tight_layout(); return save_chart(fig,'fd_vs_acc')

def chart_module_contrib():
    fig,ax=plt.subplots(figsize=(8,4.5))
    modules=['162-dim\nbaseline','+ WavLM\n(SSL帧级)','+ Adapter\n(分布校准)','+ Prosody\nPooling','+ Full\nFine-tune']
    accs, deltas, colors = [35.0,78.34,79.51,81.24,81.96],[0,43.34,1.17,1.73,0.72],['#999','#2196F3','#FF9800','#4CAF50','#9C27B0']
    bars=ax.bar(modules,accs,color=colors,width=0.55,edgecolor='white')
    for bar,acc,delta in zip(bars,accs,deltas):
        ax.text(bar.get_x()+bar.get_width()/2.,bar.get_height()+1.2,f'{acc:.1f}%',ha='center',fontsize=11,fontweight='bold',color='#333')
        if delta>1: ax.text(bar.get_x()+bar.get_width()/2.,bar.get_height()-5,f'+{delta:.1f}pp',ha='center',fontsize=9,color='white',fontweight='bold')
    ax.set_ylim(0,90); ax.set_ylabel('Test Accuracy (%)'); ax.set_title('模块贡献分解')
    fig.tight_layout(); return save_chart(fig,'module_contrib')

def chart_phase3_ablation():
    fig,ax=plt.subplots(figsize=(9,4.5))
    x,w=range(5),0.25
    seeds,key_map=['seed_42','seed_123','seed_456'],['A1_baseline','A2_stat_prior','A2b_rand_init','B3_prosody','A3_full_finetune']
    for j,(seed,c) in enumerate(zip(seeds,['#90CAF9','#64B5F6','#42A5F5'])):
        ax.bar([xi+j*w for xi in x],[phase3[k][seed]['test']*100 for k in key_map],w,color=c,label=f'seed {j+1}',edgecolor='white')
    for i,k in enumerate(key_map):
        ax.text(i+w,test_means[i]*100+1,f'{test_means[i]*100:.1f}%',ha='center',fontsize=10,fontweight='bold')
    ax.set_xticks([xi+w for xi in x]); ax.set_xticklabels(experiments,fontsize=9)
    ax.set_ylabel('Test Accuracy (%)'); ax.set_title('Phase 3 全消融 (3-way split, 3 seeds)')
    ax.legend(loc='lower right',fontsize=8); ax.set_ylim(70,90)
    fig.tight_layout(); return save_chart(fig,'phase3_ablation')

def chart_old_baselines():
    fig,ax=plt.subplots(figsize=(7,4))
    x,w=range(4),0.3
    ax.bar([xi-w/2 for xi in x],old_before,w,color='#FF4136',label='修复前(泄露)',edgecolor='white')
    ax.bar([xi+w/2 for xi in x],old_after,w,color='#4CAF50',label='修复后(真实)',edgecolor='white')
    for i in x:
        ax.annotate('',xy=(i+w/2,old_after[i]),xytext=(i-w/2,old_before[i]),arrowprops=dict(arrowstyle='->',color='#333',lw=1.5))
    ax.set_xticks(x); ax.set_xticklabels(old_models); ax.set_ylabel('Accuracy (%)')
    ax.set_title('数据泄露修复的影响'); ax.legend(fontsize=9); ax.set_ylim(0,92)
    fig.tight_layout(); return save_chart(fig,'old_baselines')

def chart_phase1_ssl():
    fig,ax=plt.subplots(figsize=(6,4))
    models=['162-dim\nmean-pooled','emotion2vec\nBase','WavLM\nBase']
    accs,colors=[35,66,80.66],['#999','#FF9800','#2196F3']
    bars=ax.bar(models,accs,color=colors,width=0.5,edgecolor='white')
    for bar,acc in zip(bars,accs): ax.text(bar.get_x()+bar.get_width()/2.,bar.get_height()+1.5,f'{acc}%',ha='center',fontsize=13,fontweight='bold')
    ax.set_ylabel('Validation Accuracy (%)'); ax.set_title('Phase 1: SSL 基线验证')
    ax.set_ylim(0,90); fig.tight_layout(); return save_chart(fig,'phase1_ssl')

def chart_3phase():
    fig,ax=plt.subplots(figsize=(7,4.2))
    phases=['原始论文\n(2025)','泄露修复后\n(4月25日)','当前方案\n(分布驱动)']
    accs,colors=[86.82,35,81.96],['#FF4136','#999','#2196F3']
    bars=ax.bar(phases,accs,color=colors,width=0.5,edgecolor='white')
    for bar,acc in zip(bars,accs): ax.text(bar.get_x()+bar.get_width()/2.,bar.get_height()+1,f'{acc}%',ha='center',fontsize=14,fontweight='bold')
    ax.axhline(y=16.7,color='#999',linestyle='--'); ax.text(2.5,17.5,'随机基线 16.7%',ha='center',fontsize=9,color='#999')
    ax.set_ylim(0,95); ax.set_ylabel('Accuracy (%)'); ax.set_title('三阶段演进')
    fig.tight_layout(); return save_chart(fig,'3phase')

def chart_cross_lang():
    fig,ax=plt.subplots(figsize=(6,4))
    exps=['同语言混合\n(MY)','英语→泰卢固语\n(English→Telugu)','泰卢固语→英语\n(Telugu→English)']
    accs,colors=[81.24,19.68,28.15],['#4CAF50','#FF4136','#FF9800']
    bars=ax.bar(exps,accs,color=colors,width=0.5,edgecolor='white')
    for bar,acc in zip(bars,accs): ax.text(bar.get_x()+bar.get_width()/2.,bar.get_height()+1,f'{acc:.1f}%',ha='center',fontsize=12,fontweight='bold')
    ax.axhline(y=16.7,color='#999',linestyle='--'); ax.text(2.1,17.5,'随机',ha='center',fontsize=9,color='#999')
    ax.set_ylabel('Test Accuracy (%)'); ax.set_title('跨语言迁移：语言的"分布鸿沟"')
    ax.set_ylim(0,90); fig.tight_layout(); return save_chart(fig,'cross_lang')

# Generate all charts
print("Generating charts...")
charts = {
    'fd_vs_acc': chart_fd_vs_acc(), 'module_contrib': chart_module_contrib(),
    'phase3_ablation': chart_phase3_ablation(), 'old_baselines': chart_old_baselines(),
    'phase1_ssl': chart_phase1_ssl(), '3phase': chart_3phase(), 'cross_lang': chart_cross_lang(),
}
print("Charts done.")

# ═══════════════════════ PPT Builder ═══════════════════════
prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
W, H = prs.slide_width, prs.slide_height

def blank_slide(): return prs.slides.add_slide(prs.slide_layouts[6])
def dark_bg(s):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = DARK_BG
def light_bg(s):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
def section_bar(s,label=None,color=PRIMARY):
    bar=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,W,Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb=color; bar.line.fill.background()
    if label:
        add_text(s,Inches(0.8),Inches(0.12),Inches(11),Inches(0.3),label,fs=10,color=GRAY)

def add_box(slide,l,t,w,h,fill=None,border=None):
    shape=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h)
    shape.fill.background()
    if fill: shape.fill.solid(); shape.fill.fore_color.rgb=_to_rgb(fill)
    if border: shape.line.color.rgb=_to_rgb(border); shape.line.width=Pt(1)
    else: shape.line.fill.background()
    return shape

def add_text(slide,l,t,w,h,text,fs=18,color=BLACK,bold=False,align=PP_ALIGN.LEFT,fn='Microsoft YaHei'):
    tb=slide.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=text; p.font.size=Pt(fs); p.font.color.rgb=color
    p.font.bold=bold; p.font.name=fn; p.alignment=align
    return tb

def add_ml(slide,l,t,w,h,lines,color=BLACK,fn='Microsoft YaHei'):
    """lines: list of (text, bold, font_size)"""
    tb=slide.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True
    for i,(txt,bld,fs) in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text=txt; p.font.size=Pt(fs); p.font.color.rgb=color; p.font.bold=bld; p.font.name=fn
    return tb

def add_pn(slide,num,total=25):
    add_text(slide,W-Inches(1.5),H-Inches(0.5),Inches(1.2),Inches(0.4),f'{num}/{total}',fs=10,color=GRAY,align=PP_ALIGN.RIGHT)

def add_img(slide,path,l,t,w,h=None):
    if h is None: h=w*0.5625
    if os.path.exists(path): return slide.shapes.add_picture(path,l,t,w,h)
    else:
        bx=add_box(slide,l,t,w,h,fill=LIGHT_GRAY,border=GRAY)
        add_text(slide,l+Inches(0.2),t+h/2-Inches(0.2),w-Inches(0.4),Inches(0.4),'[图片]',fs=11,color=GRAY,align=PP_ALIGN.CENTER); return bx

def add_code_block(slide,l,t,w,h,code_text,title=None,fs=10):
    """Add a code-looking block with dark background."""
    box=add_box(slide,l,t,w,h,fill=RGBColor(0x2D,0x2D,0x2D))
    if title:
        add_text(slide,l+Inches(0.15),t+Inches(0.05),w-Inches(0.3),Inches(0.25),title,fs=9,color=GRAY,bold=True)
    add_text(slide,l+Inches(0.2),t+(Inches(0.3) if title else Inches(0.1)),w-Inches(0.4),h-Inches(0.4),code_text,fs=fs,color=RGBColor(0xE0,0xE0,0xE0),fn='Consolas')
    return box

def add_table(slide,l,t,col_widths,headers,rows,fs=12):
    n_rows,n_cols=len(rows)+1,len(headers)
    shape=slide.shapes.add_table(n_rows,n_cols,l,t,sum(col_widths),Inches(0.38)*n_rows)
    tbl=shape.table
    for j,(h,w_) in enumerate(zip(headers,col_widths)): tbl.columns[j].width=w_
    for j,h in enumerate(headers):
        c=tbl.cell(0,j); c.text=h
        c.fill.solid(); c.fill.fore_color.rgb=PRIMARY
        for p in c.text_frame.paragraphs: p.font.size=Pt(fs); p.font.color.rgb=WHITE; p.font.bold=True; p.font.name='Microsoft YaHei'; p.alignment=PP_ALIGN.CENTER
    for i,row in enumerate(rows):
        for j,val in enumerate(row):
            c=tbl.cell(i+1,j); c.text=str(val)
            if i%2==0: c.fill.solid(); c.fill.fore_color.rgb=LIGHT_GRAY
            for p in c.text_frame.paragraphs: p.font.size=Pt(fs); p.font.color.rgb=BLACK; p.font.name='Microsoft YaHei'; p.alignment=PP_ALIGN.CENTER
    return shape

# ═══════════════════════════════════════════
# SLIDE 1 — Cover
# ═══════════════════════════════════════════
s=blank_slide(); dark_bg(s)
add_text(s,Inches(2),Inches(1.5),Inches(9),Inches(1.5),'儿童语音情绪识别',fs=48,color=WHITE,bold=True,align=PP_ALIGN.CENTER)
add_text(s,Inches(2),Inches(2.8),Inches(9),Inches(1),'从数据泄露到分布驱动框架',fs=32,color=PRIMARY,align=PP_ALIGN.CENTER)
line=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(4.5),Inches(3.9),Inches(4),Inches(0.03))
line.fill.solid(); line.fill.fore_color.rgb=PRIMARY; line.line.fill.background()
add_text(s,Inches(2),Inches(4.3),Inches(9),Inches(0.6),'项目进展汇报',fs=24,color=GRAY,align=PP_ALIGN.CENTER)
add_text(s,Inches(2),Inches(5.3),Inches(9),Inches(0.5),'2026年4月29日',fs=16,color=GRAY,align=PP_ALIGN.CENTER)
add_text(s,Inches(2),Inches(5.8),Inches(9),Inches(0.5),'GitHub: Xuzc317/child-speech-emotion-recognition',fs=12,color=GRAY,align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# SLIDE 2 — Roadmap (REDESIGNED: timeline style)
# ═══════════════════════════════════════════
s=blank_slide(); light_bg(s); section_bar(s)
add_text(s,Inches(0.8),Inches(0.3),Inches(11),Inches(0.8),'汇报路线图',fs=36,color=BLACK,bold=True)

# Timeline bar
bar_y = Inches(2.8)
line=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(1),bar_y,Inches(11.3),Inches(0.04))
line.fill.solid(); line.fill.fore_color.rgb=PRIMARY; line.line.fill.background()

# 5 timeline nodes
nodes = [
    ('2025年\n5-7月','原始论文\n86.82%','#FF4136'),
    ('2026年\n4月25-26日','发现泄露\n真实基线~35%','#FF9800'),
    ('4月26日','系统性修复\n重评估','#2196F3'),
    ('4月27-28日','SSL+三模块\n分布驱动框架','#4CAF50'),
    ('4月28-29日','整合实验\n81.96%','#9C27B0'),
]
for i,(date,label,color) in enumerate(nodes):
    cx = Inches(1.5 + i*2.4)
    # Dot on line
    dot=s.shapes.add_shape(MSO_SHAPE.OVAL,cx+Inches(0.7),bar_y-Inches(0.08),Inches(0.2),Inches(0.2))
    dot.fill.solid(); dot.fill.fore_color.rgb=_to_rgb(color); dot.line.fill.background()
    # Date above
    add_text(s,cx,bar_y-Inches(0.7),Inches(1.6),Inches(0.55),date,fs=11,color=GRAY,align=PP_ALIGN.CENTER)
    # Label below (in a card)
    card=add_box(s,cx,bar_y+Inches(0.4),Inches(1.6),Inches(0.8),fill=color)
    add_text(s,cx+Inches(0.05),bar_y+Inches(0.5),Inches(1.5),Inches(0.6),label,fs=12,color=WHITE,bold=True,align=PP_ALIGN.CENTER)

# Key message
add_ml(s,Inches(0.8),Inches(4.6),Inches(11),Inches(2.5),[
    ('三次推进，两次来自外部反馈',True,22),
    ('',False,6),
    ('• 面试时 CTO 质疑增强策略 → 发现 63% speaker 重叠的数据泄露根因',False,16),
    ('• 入职时 CTO 纠正我对自监督学习（SSL）的判断 → 找到分布驱动新方向',False,16),
    ('• 自主推进 Phase 1→4：SSL 选型 → 三模块实现 → 云端全消融 → 跨语言实验',False,16),
    ('',False,6),
    ('核心叙事：不是"更强的模型"，而是从儿童语音的统计分布出发，重新思考 SER 每个环节应该怎么做',False,16),
])
add_pn(s,2)

# ═══════════════════════════════════════════
# SLIDE 3 — Original Plan Overview
# ═══════════════════════════════════════════
s=blank_slide(); light_bg(s); section_bar(s,label='第一部分：原始方案回顾')
add_text(s,Inches(0.8),Inches(0.3),Inches(11),Inches(0.8),'原始论文方案（2025年5-7月）',fs=36,color=BLACK,bold=True)

add_ml(s,Inches(0.8),Inches(1.4),Inches(5.5),Inches(3.5),[
    ('方案设计',True,22),('',False,6),
    ('• 任务：BESD 儿童双语数据集，6 类情绪分类',False,14),
    ('• 特征：手工声学特征 → mean pooling 压缩为固定向量',False,14),
    ('• 增强：每 WAV → 3x（原始 + 加噪 + 拉伸变调）',False,14),
    ('• 模型：DrseCNN（ResNet 残差 + SE 注意力 + 多阶段卷积）',False,14),
    ('• 参数量：~11.37M',False,14),
],color=BLACK)

add_table(s,Inches(6.8),Inches(1.4),[Inches(2.2),Inches(1.5),Inches(1.5)],
    ['模型','Accuracy','参数量'],
    [['DrseCNN','86.82%','11.37M'],['CNN-BiLSTM','72.26%','-'],['CNN','72.14%','-'],['Transformer','69.48%','-']])
add_text(s,Inches(6.8),Inches(4.2),Inches(5),Inches(0.5),'看起来效果很好，消融实验完整',fs=13,color=GRAY,align=PP_ALIGN.CENTER)
add_pn(s,3)

# ═══════════════════════════════════════════
# SLIDE 4 — Feature extraction details (NEW: explain 94→162)
# ═══════════════════════════════════════════
s=blank_slide(); light_bg(s); section_bar(s,'')
add_text(s,Inches(0.8),Inches(0.3),Inches(11),Inches(0.8),'原始方案的特征提取：94维是怎么变成162维的',fs=34,color=BLACK,bold=True)

# Left: feature composition
add_ml(s,Inches(0.8),Inches(1.3),Inches(5.8),Inches(3),[
    ('第一步：对每一条WAV，提取5种手工声学特征',True,18),('',False,4),
    ('  ZCR(过零率) 1维 + Chroma(色度) 12维 + MFCC(梅尔倒谱) 40维 + RMS(能量) 1维 + Mel(梅尔频谱) 40维 = 共94维',False,14),
    ('  → ZCR反映过零速率，Chroma映射12半音阶，MFCC模拟人耳低频感知，RMS反映响度，Mel是频域能量分布',False,13),
    ('',False,4),
],color=BLACK)

# Right: what happens with augmentation
add_ml(s,Inches(7.2),Inches(1.3),Inches(5.5),Inches(3),[
    ('第二步：3x增强把94维变成可训练矩阵',True,18),('',False,4),
    ('每一条WAV做了3个版本：',False,15),
    (' ① 原始版本  → extract_features() → 94维向量',False,14),
    (' ② 加噪版本  → noise() + extract_features() → 94维',False,14),
    (' ③ 拉伸+变调 → stretch()+pitch()+extract_features() → 94维',False,14),
    ('',False,6),
    ('三个版本垂直堆叠 → (3, 94) 的矩阵',False,16),
    ('→ 这就是"162维"的来源（论文写成162维，实际是3×54或3×94）',False,14),
    ('→ mean() 压平 → (94,) 或 (162,) 固定向量',False,14),
],color=BLACK)

# Bottom: the critical problem
add_box(s,Inches(0.8),Inches(4.8),Inches(11.5),Inches(2.3),fill=RGBColor(0xFF,0xF0,0xF0),border=RED)
add_ml(s,Inches(1.2),Inches(4.9),Inches(10.8),Inches(2),[
    ('关键问题：两个层次的信息丢失',True,20),('',False,4),
    ('1. 时间维度丢失：对每一条WAV的T帧特征做了 np.mean(所有帧) → 只得到1个值代表整条语音',False,15),
    ('  例：ZCR = np.mean(librosa.feature.zero_crossing_rate(y=data).T, axis=0)  ← 所有200帧过零率被压成1个数',False,14),
    ('2. 特征维度混淆：最终 (3,94) 矩阵 → Conv1d 的 kernel 在"94维特征拼接轴"上滑动，而非"3个增强版本的时间轴"',False,14),
    ('',False,4),
    ('→ 所有"时序模型"（LSTM/Transformer/BiLSTM）都在没有时间维度的数据上运行。这是结构性错误。',True,16),
],color=BLACK)
add_pn(s,4)

# ═══════════════════════════════════════════
# SLIDE 5 — CTO Interview + Data Leak (WITH CODE DETAILS)
# ═══════════════════════════════════════════
s=blank_slide(); light_bg(s); section_bar(s,'')
add_text(s,Inches(0.8),Inches(0.3),Inches(11),Inches(0.8),'转折一：CTO的一个问题，让我重新审视整个实验',fs=36,color=RED,bold=True)

qbox=add_box(s,Inches(0.8),Inches(1.2),Inches(11.5),Inches(1.0),fill=RGBColor(0xFF,0xF0,0xF0),border=RED)
add_text(s,Inches(1.2),Inches(1.35),Inches(10.8),Inches(0.75),
    'CTO："你先对每条语音做了加噪、拉伸、变调，生成了三个增强版本，然后再对整个数据集做随机划分。'
    '你有没有想过，同一条原始语音的增强版本，可能同时出现在训练集和测试集里？"',
    fs=16,color=RED)

# Two-column code comparison
add_text(s,Inches(0.8),Inches(2.4),Inches(5.5),Inches(0.5),'原始代码的执行顺序（问题所在）：',fs=16,color=BLACK,bold=True)
add_code_block(s,Inches(0.8),Inches(2.8),Inches(5.5),Inches(2.2),
    '# 步骤1：先对每条WAV做3x增强\n'
    'for wav in all_wavs:\n'
    '    feat_raw  = extract_features(wav)      # 94维\n'
    '    feat_noisy = extract_features(noise(wav))  # 94维\n'
    '    feat_pitch = extract_features(pitch(stretch(wav))) # 94维\n'
    '    result = np.vstack([feat_raw, feat_noisy, feat_pitch])\n'
    '    all_data.append(result)  # (3,94) per wav\n'
    '\n'
    '# 步骤2：数据全部生成后，再随机洗牌切分\n'
    'np.random.seed(42)\n'
    'indices = np.random.permutation(12537)\n'
    'split = int(0.7 * 12537)  # 8775 / 3762',
    title='englishdata.py (原始代码)')

add_text(s,Inches(7),Inches(2.4),Inches(5.5),Inches(0.5),'应该的正确顺序：',fs=16,color=GREEN,bold=True)
add_code_block(s,Inches(7),Inches(2.8),Inches(5.5),Inches(2.2),
    '# 步骤1：先按说话人划分训练/测试集\n'
    'train_speakers, test_speakers = split_speakers(...)\n'
    '# → 165人训练 / 72人测试，0重叠\n'
    '\n'
    '# 步骤2：训练集和测试集各自独立增强\n'
    'train_data = []\n'
    'for wav in train_wavs:\n'
    '    # 3x增强 → 仅加入训练数据\n'
    '    train_data.append(augment(wav))\n'
    '\n'
    'test_data = []\n'
    'for wav in test_wavs:\n'
    '    # 测试集不增强，或用独立参数增强\n'
    '    test_data.append(extract_features(wav))',
    title='正确的流程')

# Bottom explanation
add_ml(s,Inches(0.8),Inches(5.3),Inches(11.5),Inches(1.8),[
    ('为什么 np.random.permutation 本身没问题，但加上增强顺序就错了？',True,16),('',False,4),
    ('同一条WAV的3个增强版本高度相似（加了同一条噪声种子）。random.shuffle 无法识别它们来自同一个人——',False,14),
    ('因为代码里根本没有告诉 shuffle "这三个样本来自同一个说话人"。',False,14),
    ('验证方式：从文件名提取 speaker ID，检查训练/测试集中同一 speaker 的出现比例 → 63%。',False,14),
],color=BLACK)
add_pn(s,5)

# ═══════════════════════════════════════════
# SLIDE 6 — More Problems (FIXED: #4 changed to sample count discrepancy)
# ═══════════════════════════════════════════
s=blank_slide(); light_bg(s); section_bar(s,'')
add_text(s,Inches(0.8),Inches(0.3),Inches(11),Inches(0.8),'顺着这个问题，挖出了更多隐患',fs=36,color=BLACK,bold=True)

add_table(s,Inches(0.5),Inches(1.4),
    [Inches(0.4),Inches(5.5),Inches(6.5),Inches(1)],
    ['#','问题','详细说明','严重度'],
    [['1','数据泄露（speaker重叠）','增强后随机划分→同一说话人的3个增强版本有63%概率出现在训练+测试集','🔴'],
     ['2','时间维度完全丢失','mean pooling 抹平所有帧→Conv1d 在ZCR/Chroma/MFCC的特征拼接维上滑动→不是真正的时间建模','🔴'],
     ['3','增强顺序错误','先增强后划分→测试集中混入增强数据→测试分布被人为扭曲','🟠'],
     ['4','样本条数不一致','原始BESD MY共4,179条WAV，3x增强应=12,537条，但论文统计/实验中出现了不一致的数据量，导致描述混乱','🟠'],
     ['5-9','其他问题','特征维度标注模糊（94 vs 162）、网络层数硬编码、消融实验描述不一致等','🟡']])

add_ml(s,Inches(0.8),Inches(4.6),Inches(11),Inches(2.5),[
    ('最隐蔽的问题——时间维度丢失（为什么这是个结构性错误）：',True,18),
    ('',False,4),
    ('WAV(2.5s) → librosa 帧级特征 (200帧 × 94维) → np.mean(所有帧, axis=0) → (94,) 向量',False,15),
    ('                                           ↑ 200帧的时间信息在这一步完全消失',False,14),
    ('',False,4),
    ('于是：Conv1d kernel → 在"ZCR+Chroma+MFCC+RMS+Mel"的特征拼接轴上滑动，不是时间轴',False,14),
    ('     Transformer → (B, 1, 94) → self-attention 在 T=1 上做，变成恒等映射',False,14),
    ('     CRNN LSTM → 把下采样后的 ~40 个"位置"当时间步 → 位置 i→i+1 是不同特征组的边界',False,14),
])
add_pn(s,6)

# ═══════════════════════════════════════════
# SLIDE 7 — Systematic Fix
# ═══════════════════════════════════════════
s=blank_slide(); light_bg(s); section_bar(s,label='第二部分：修复与真实基线')
add_text(s,Inches(0.8),Inches(0.3),Inches(11),Inches(0.8),'系统性修复（2026年4月25-26日）',fs=36,color=BLACK,bold=True)

add_table(s,Inches(0.8),Inches(1.4),[Inches(3),Inches(8.5)],
    ['修复项','操作'],
    [['说话人划分','从文件名提取 speaker ID → normalize_speaker_id() 修正变体标注'],
     ['分层划分','Profile-based stratification（按情感组合分组后 70/30 切分训练/测试）'],
     ['数据集扩展','ENGLISH → MY（ENGLISH + TELUGU，共 4,179 WAVs，237 个说话人）'],
     ['增强顺序','先划分 → 再增强（训练集和测试集各自独立处理增强流程）'],
     ['旧权重管理','全部移至 legacy/ 目录，标注无效，防止误用']])

add_text(s,Inches(0.8),Inches(4.2),Inches(11),Inches(0.5),'修复后的数据状态：',fs=20,color=BLACK,bold=True)
add_table(s,Inches(0.8),Inches(4.8),[Inches(4),Inches(4),Inches(4)],
    ['指标','训练集','测试集'],
    [['说话人数','165','72（0重叠 ✅）'],
     ['样本数（增强后）','~8,700','~3,837'],
     ['每类样本数','~1,400-1,470','~630-645']])
add_pn(s,7)

# ═══════════════════════════════════════════
# SLIDE 8 — Real Baseline After Fix
# ═══════════════════════════════════════════
s=blank_slide(); light_bg(s); section_bar(s,'')
add_text(s,Inches(0.8),Inches(0.3),Inches(11),Inches(0.8),'数据泄露修复后——真实基线',fs=36,color=RED,bold=True)

add_table(s,Inches(0.8),Inches(1.3),[Inches(3),Inches(2.5),Inches(2.5),Inches(2.5)],
    ['模型','修复前','修复后','变化'],
    [['DrseCNN','85.99%','~35%','-51pp'],['CNN','72.14%','~30%','-42pp'],
     ['BiLSTM','72.26%','~30%','-42pp'],['Transformer','69.48%','~28%','-41pp']])
add_img(s,charts['old_baselines'],Inches(0.3),Inches(4.2),Inches(5.8))

add_ml(s,Inches(6.8),Inches(4.2),Inches(5.5),Inches(2.8),[
    ('三个关键发现：',True,18),('',False,4),
    ('1. 所有模型一起崩到 30-40%（随机基线 16.7%）→ 证明问题不是"模型不够好"',False,14),
    ('2. 模型间"性能梯度"消失 → 原来的排序（DrseCNN>BiLSTM>CNN>Transformer）本质上是泄露程度的排序',False,14),
    ('3. 162-dim mean-pooled 的信息上限 ≈ 35% → 特征表示才是瓶颈，不是分类器',False,14),
    ('',False,6),
    ('这是整个项目最重要的转折点。',True,18),
    ('→ 结论：不是模型的问题，是特征从根本上就错了。',True,16),
],color=BLACK)
add_pn(s,8)

# ═══════════════════════════════════════════
# SLIDE 9 — Phase Overview (NEW: before Phase 1)
# ═══════════════════════════════════════════
s=blank_slide(); light_bg(s); section_bar(s,label='第三部分：新方案设计')
add_text(s,Inches(0.8),Inches(0.3),Inches(11),Inches(0.8),'新方案研发：四个阶段的完整路线',fs=36,color=BLACK,bold=True)

phases = [
    ('Phase 1\n基线验证','验证 SSL 帧级特征\n是否能突破 35% 上限','1天\n4月26日','✅\n80.66%','#2196F3'),
    ('Phase 2\n三模块实现','逐个实现：Adapter\n+ Pooling + Augmentation','2天\n4月27-28日','✅\n86.42%','#4CAF50'),
    ('Phase 3\n整合实验','云端全消融\n3-way split + 3 seeds','1天\n4月28日','✅\n81.96%','#9C27B0'),
    ('Phase 4\n论文撰写','补充实验（跨语言）\n+ 论文初稿 v3','进行中','🔄\n待完成','#FF9800'),
]
for i,(title,desc,time,result,color) in enumerate(phases):
    cx=Inches(0.6+i*3.1)
    add_box(s,cx,Inches(1.5),Inches(2.8),Inches(3.8),fill=color)
    add_text(s,cx+Inches(0.1),Inches(1.6),Inches(2.6),Inches(0.8),title,fs=22,color=WHITE,bold=True,align=PP_ALIGN.CENTER)
    add_text(s,cx+Inches(0.1),Inches(2.4),Inches(2.6),Inches(1.0),desc,fs=13,color=WHITE,align=PP_ALIGN.CENTER)
    add_text(s,cx+Inches(0.1),Inches(3.4),Inches(2.6),Inches(0.5),time,fs=11,color=WHITE,align=PP_ALIGN.CENTER)
    add_text(s,cx+Inches(0.1),Inches(4.0),Inches(2.6),Inches(0.8),result,fs=20,color=WHITE,bold=True,align=PP_ALIGN.CENTER)

add_ml(s,Inches(0.8),Inches(5.8),Inches(11),Inches(1.5),[
    ('每个阶段的进入条件和验收条件：',True,16),
    ('• Phase 1 进入条件：确认 BESD MY 可用、CUDA 环境正常',False,14),
    ('• Phase 2 进入条件：Phase 1 验证 SSL 显著优于 35% → 进入条件满足（80.66% > 35%）',False,14),
    ('• Phase 3 进入条件：Phase 2 三个模块均有正向消融结果 → 上云端跑完整实验',False,14),
    ('• Phase 4 进入条件：Phase 3 结果稳定、数据无泄露 → 开始撰写论文',False,14),
])
add_pn(s,9)

# ═══════════════════════════════════════════
# SLIDE 10 — CTO Interview 2: SSL
# ═══════════════════════════════════════════
s=blank_slide(); light_bg(s); section_bar(s,'')
add_text(s,Inches(0.8),Inches(0.3),Inches(11),Inches(0.8),'转折二：在我差点放弃自监督学习的时候',fs=36,color=GREEN,bold=True)

add_ml(s,Inches(0.8),Inches(1.2),Inches(11.5),Inches(1.0),[
    ('入职时跟CTO聊起项目：我说修复泄露后准确率只有35%，试过wav2vec2做自监督特征提取但效果不好，所以觉得自监督学习跟这个项目不匹配。',False,16),
],color=BLACK)

qbox=add_box(s,Inches(0.8),Inches(2.0),Inches(11.5),Inches(1.3),fill=RGBColor(0xE8,0xF5,0xE9),border=GREEN)
add_ml(s,Inches(1.2),Inches(2.1),Inches(10.8),Inches(1.1),[
    ('CTO的回应：',True,18),
    ('"wav2vec2 效果不好，不代表自监督学习这条路不对。wav2vec2 是 2020 年的模型，预训练数据是成人朗读。现在有更新的模型，比如 WavLM 在语音任务上普遍比 wav2vec2 强，emotion2vec 更是专门为情绪识别设计的。"',False,15),
],color=BLACK)

add_ml(s,Inches(0.8),Inches(3.6),Inches(11.5),Inches(3.5),[
    ('他从原理层面梳理了关键点：',True,18),('',False,4),
    ('• 自监督预训练模型学到的是帧级别的连续表征 → 天然保留时间结构 → 正好对症下药',False,15),
    ('• 我之前 162-dim mean-pooled 的问题恰恰是丢掉了时间维度 → SSL 正好解决了这个根本矛盾',False,15),
    ('• 不是自监督不适合，而是需要选对模型 + 根据儿童语音特点做适配',False,15),
    ('',False,8),
    ('由此形成了"分布驱动三模块"的整体思路：',True,18),
    ('',False,4),
    ('直接拿成人预训练的 SSL 模型还不够——儿童语音在音高范围、语速、共振峰等方面与成人有系统性的分布差异。',False,15),
    ('→ 需要三个技术模块来弥合这个分布偏移：分布校准适配器 / 韵律引导池化 / 分布约束增强',False,15),
])
add_pn(s,10)

# ═══════════════════════════════════════════
# SLIDE 11 — What is SSL frame-level feature? (NEW)
# ═══════════════════════════════════════════
s=blank_slide(); light_bg(s); section_bar(s,'')
add_text(s,Inches(0.8),Inches(0.3),Inches(11),Inches(0.8),'关键概念：什么是"SSL 帧级特征"？',fs=36,color=BLACK,bold=True)

# Left: old way
add_box(s,Inches(0.5),Inches(1.3),Inches(5.8),Inches(3.0),fill=LIGHT_GRAY,border=GRAY)
add_text(s,Inches(0.8),Inches(1.4),Inches(5.3),Inches(0.5),'旧方案：手工特征 + mean pooling',fs=20,color=RED,bold=True)
add_ml(s,Inches(0.8),Inches(1.9),Inches(5.3),Inches(2.2),[
    ('WAV语音 (~2.5秒, 16kHz采样 = 40,000个采样点)',False,14),
    ('    ↓ 分帧：每25ms一帧，帧移10ms',False,14),
    ('~200帧 × 每帧提取94个手工特征值',False,14),
    ('    ↓ np.mean(所有200帧, axis=0)',False,14),
    ('一个 (94,) 向量 → 整条语音被压缩成94个数',False,15),
    ('    ↓',False,12),
    ('时间信息：✗ 完全丢失（200帧 → 1个均值）',True,15),
    ('信息量：94个数值 / 条语音',False,14),
],color=BLACK)

# Right: new way
add_box(s,Inches(7),Inches(1.3),Inches(5.8),Inches(3.0),fill=RGBColor(0xE3,0xF2,0xFD),border=PRIMARY)
add_text(s,Inches(7.3),Inches(1.4),Inches(5.3),Inches(0.5),'新方案：SSL 帧级特征',fs=20,color=PRIMARY,bold=True)
add_ml(s,Inches(7.3),Inches(1.9),Inches(5.3),Inches(2.2),[
    ('WAV语音 (~2.5秒, 16kHz采样 = 40,000个采样点)',False,14),
    ('    ↓ WavLM 模型（94M参数，在94,000小时语音上预训练）',False,14),
    ('每20ms输出一帧，每帧768维 → ~200帧',False,14),
    ('    ↓ 保留所有帧，不压缩',False,14),
    ('一个 (200, 768) 矩阵 → 每个时间步都有独立表示',False,15),
    ('    ↓',False,12),
    ('时间信息：✓ 完整保留（200帧 → 200个768维向量）',True,15),
    ('信息量：153,600个数值 / 条语音（1600倍）',False,14),
],color=BLACK)

# Bottom: why SSL works
add_ml(s,Inches(0.8),Inches(4.6),Inches(11.5),Inches(2.5),[
    ('为什么"预训练"+"帧级"这两个特性是关键？',True,20),('',False,4),
    ('1. "预训练"：WavLM 在94,000小时的LibriSpeech+GigaSpeech等成人语音上做过自监督学习——它已经学会了如何从原始波形中提取有意义的声学模式，不需要从零训练。类似 GPT 已经读过海量文本，WavLM 已经"听过"海量语音。',
     False,14),
    ('2. "帧级"：预训练阶段的设计目标是预测被掩码的帧——所以它必须为每个时间帧生成有意义的表示。这意味着输出天然保留了时间结构（每20ms一个768维向量），不像 mean pooling 那样把200帧压成1个数。',
     False,14),
    ('3. 类比：旧方案 = 把一部90分钟电影压成一张剧照然后分析情节；新方案 = 保留完整90分钟画面和声音，让模型自己判断哪个时刻的情绪最重要。',False,14),
])
add_pn(s,11)

# ═══════════════════════════════════════════
# SLIDE 12 — Phase 1 Results
# ═══════════════════════════════════════════
s=blank_slide(); light_bg(s); section_bar(s,'')
add_text(s,Inches(0.8),Inches(0.3),Inches(11),Inches(0.8),'Phase 1 — 验证"保留时间维度"的价值',fs=36,color=BLACK,bold=True)

add_img(s,charts['phase1_ssl'],Inches(0.3),Inches(1.2),Inches(5.8))
add_table(s,Inches(6.8),Inches(1.2),[Inches(2.5),Inches(1.5),Inches(2)],
    ['模型','Val Acc','备注'],
    [['旧162-dim baseline','~35%','不可用'],
     ['emotion2vec Base','~66%','震荡较大'],
     ['WavLM Base','80.66%','稳定，方差小']])

add_ml(s,Inches(6.8),Inches(3.3),Inches(5.5),Inches(3.5),[
    ('实验方法：Frozen WavLM 提取帧级特征 (N,200,768) → 仅训练一个简单的线性分类器（Linear Probe）',True,14),('',False,4),
    ('结论：',True,20),('',False,4),
    ('SSL 帧级特征 + 保留时间维',False,18),
    ('→ 35% → 80.66%，+45pp',False,18),
    ('→ 确认方向正确，以 WavLM 为主模型',False,18),('',False,8),
    ('备注：emotion2vec 后续在 Phase 3 发现其 plus_large 版本仅 22.45%——预处理语言不匹配，WavLM 通用性更稳',
     False,12),
],color=BLACK)
add_pn(s,12)

# ═══════════════════════════════════════════
# SLIDE 13 — Module 1: Adapter (IMPROVED: clear what it is)
# ═══════════════════════════════════════════
s=blank_slide(); light_bg(s); section_bar(s,'')
add_text(s,Inches(0.8),Inches(0.3),Inches(11),Inches(0.8),'模块1 — AcousticCalibrationAdapter',fs=36,color=BLACK,bold=True)

# Left: plain language
add_ml(s,Inches(0.8),Inches(1.3),Inches(6),Inches(3.2),[
    ('它到底是个什么东西？',True,22),('',False,4),
    ('一句话：在 WavLM 冻结的输出上加一个可学习的"翻译层"。',False,17),('',False,6),
    ('WavLM 在成人语音上预训练 → 它"听"儿童语音时，隐层表示会带偏移（就像北京人第一次听四川话，能听懂但不是100%准确）。',False,15),
    ('',False,4),
    ('Adapter = 一个轻量的非线性映射模块（~900K参数）：',False,15),
    ('  Input: WavLM 第k层的输出 (T, 768)',False,14),
    ('  → LayerNorm → Linear(768→192) → GELU → Linear(192→768)',False,14),
    ('  → 残差连接：Output = Input + 调整量',False,14),
    ('  Output: 校准后的表示 (T, 768)',False,14),
    ('',False,4),
    ('为什么残差连接重要？Adapter 只学"纠正量"——大部分信息仍来自 WavLM，Adapter 只补儿童特有的偏移部分。',False,14),
],color=BLACK)

# Right: results table
add_table(s,Inches(7.2),Inches(1.3),[Inches(2.5),Inches(2.5),Inches(1.5),Inches(1.5)],
    ['实验','配置','Val Acc','vs 基线'],
    [['A1','无 Adapter','80.66%','—'],
     ['A2','统计先验初始化','81.45%','+0.79pp'],
     ['A2b','随机初始化','85.40%','+4.74pp']])

add_ml(s,Inches(7.2),Inches(3.6),Inches(5.5),Inches(3.5),[
    ('反直觉发现 ⚡：统计先验 < 随机初始化',True,20),('',False,4),
    ('统计先验做法：计算成人→儿童的均值+方差偏移，用线性公式补偿。',False,14),
    ('→ 效果只提升 0.79pp，说明儿童语音的分布偏移不是简单的"平移+缩放"。',False,14),
    ('',False,4),
    ('随机初始化 → 让 Adapter 自己学习校准 → +4.74pp。',False,14),
    ('→ 证明偏移是结构性的（非线性、层相关的），需要可学习模块而非固定公式。',False,14),
    ('',False,4),
    ('这是论文的第3个核心论点。',True,16),
])
add_pn(s,13)

# ═══════════════════════════════════════════
# SLIDE 14 — Module 2: Pooling (IMPROVED)
# ═══════════════════════════════════════════
s=blank_slide(); light_bg(s); section_bar(s,'')
add_text(s,Inches(0.8),Inches(0.3),Inches(11),Inches(0.8),'模块2 — TemporalImportancePooling',fs=36,color=BLACK,bold=True)

add_ml(s,Inches(0.8),Inches(1.3),Inches(6),Inches(3.5),[
    ('它到底是个什么东西？',True,22),('',False,4),
    ('一句话：用儿童的音高（F0）和能量变化来告诉模型"听哪几帧"。',False,17),('',False,6),
    ('核心逻辑：',False,16),
    ('  • 儿童说"开心"时F0会陡然升高 → 那一刻的帧最重要',False,14),
    ('  • 儿童说"伤心"时语音能量持续低落 → 关键帧分布较均匀',False,14),
    ('  • 传统 mean pooling 对所有帧一视同仁 → 关键信息被稀释',False,14),('',False,4),
    ('  (B, T, 768) SSL特征        (B, T, 2) 韵律特征',False,14),
    ('       ↓                           ↓',False,12),
    ('  学习帧权重 α(t)    ×    调节因子 β(t)    → 加权求和 → (B, 768)',False,14),
],color=BLACK)

add_ml(s,Inches(7.2),Inches(1.3),Inches(5.5),Inches(3.5),[
    ('儿童语音的韵律特异性：',True,18),('',False,4),
    ('F0 均值 495.9 Hz（成人女声 ~200Hz），标准差 651.1 Hz。',False,14),
    ('p5=75 Hz 到 p95=2,286 Hz —— 极宽的波动范围。',False,14),
    ('同一句话里音高波动可能是成人的3-5倍。',False,14),('',False,6),
    ('结果：',True,18),('',False,4),
    ('Adapter + Prosody Pooling',False,16),
    ('→ 86.42%（+1.02pp over Adapter-only）',False,16),
    ('',False,4),
    ('韵律引导 > 纯 self-attention → 可解释性优势',False,14),
    ('F0变异大的帧获得更高权重 → 与儿童情绪表达的韵律特性一致',False,13),
],color=BLACK)
add_pn(s,14)

# ═══════════════════════════════════════════
# SLIDE 15 — Module 3: Augmentation (IMPROVED: plain language)
# ═══════════════════════════════════════════
s=blank_slide(); light_bg(s); section_bar(s,'')
add_text(s,Inches(0.8),Inches(0.3),Inches(11),Inches(0.8),'模块3 — 分布约束增强是怎么实现的？',fs=36,color=BLACK,bold=True)

# Plain language explanation
add_ml(s,Inches(0.8),Inches(1.2),Inches(6),Inches(3),[
    ('白话解释：',True,22),('',False,4),
    ('传统增强：把语音像照片一样"随意P图"——音调拉高6个半音、速度放慢到0.7倍。这对成人语音 OK，因为成人语音参数变异性小，但这些操作放在儿童语音上就"P过头了"——儿童本来音高就高、语速就不稳定，再大幅增强会破坏情感特征。',
     False,14),('',False,6),
    ('我们的做法：',False,16),
    ('1. 先用 BESD 数据统计儿童语音的真实参数范围',False,14),
    ('   → F0均值 496Hz, std 651Hz, 语速分布等',False,14),
    ('2. 基于统计分布定义"合法增强空间"',False,14),
    ('   → pitch shift ±3st（成人用±6st）',False,14),
    ('   → speed 0.85-1.15（成人用 0.7-1.3）',False,14),
    ('3. 在合法空间内增强 → 样本保持"听起来像儿童"',False,14),
],color=BLACK)

# Right: FD quantification
add_ml(s,Inches(7.2),Inches(1.2),Inches(5.5),Inches(3),[
    ('怎么定量衡量"P过头了"？',True,20),('',False,4),
    ('用 Frechet Distance（FD）：',False,16),
    ('→ 衡量增强后样本分布 vs 原始分布的距离',False,14),
    ('→ FD 越大 = 增强越"离谱" = 越偏离真实儿童语音',False,14),('',False,6),
    ('FD 计算方式（类比）：',False,15),
    ('就像比较"原图"和"P过的图"的像素分布差异。',False,14),
    ('FD=0 → 完全一致；FD越大 → 偏离越远。',False,14),
    ('我们用 WavLM 隐层特征来算这个距离。',False,14),
],color=BLACK)

add_img(s,charts['fd_vs_acc'],Inches(0.3),Inches(4.3),Inches(6.2))
add_table(s,Inches(7),Inches(4.3),[Inches(1.5),Inches(1.5),Inches(1.5),Inches(1.5)],
    ['实验','FD','Val Acc','vs 无增强'],
    [['C1 无增强','0','86.42%','基线'],
     ['C3 儿童约束','8.71','66.85%','-19.57pp'],
     ['C2 成人参数','9.87','59.51%','-26.91pp'],
     ['C4 极端参数','11.99','43.25%','-43.17pp']])
add_text(s,Inches(7),Inches(6.5),Inches(5.5),Inches(0.4),
    '核心发现：FD 与准确率严格单调负相关——分布偏移可定量预测分类退化',fs=13,color=RED,bold=True)
add_pn(s,15)

# ═══════════════════════════════════════════
# SLIDE 16 — Second Leak Fix
# ═══════════════════════════════════════════
s=blank_slide(); light_bg(s); section_bar(s,label='第四部分：整合实验与最终结果')
add_text(s,Inches(0.8),Inches(0.3),Inches(11),Inches(0.8),'Phase 3 前置——第二次数据泄露修复',fs=36,color=BLACK,bold=True)

add_table(s,Inches(0.8),Inches(1.4),[Inches(3.5),Inches(3.5),Inches(5)],
    ['泄露风险','问题','修复'],
    [['adapter_init 使用全部speaker','统计先验（μ_child, σ_child）计算时包含了测试集说话人','60/20/20 三路划分，仅 train speaker 计算统计量'],
     ['test set = val set','用 test 做 early stopping → 间接泄露标签分布','val 做 early-stop，test 严格 hold-out（训练中完全不可见）']])

add_ml(s,Inches(0.8),Inches(3.5),Inches(11.5),Inches(3.5),[
    ('新增代码与验证：',True,18),
    ('• split_speakers_3way() — profile-stratified 三路划分（按情感分布分层 + 说话人互斥）',False,15),
    ('• 三层断言保证互斥：train ∩ val = ∅, train ∩ test = ∅, val ∩ test = ∅',False,15),
    ('• 每个 speaker 的所有WAV只能出现在一个集合里——这是刚性约束',False,15),('',False,6),
    ('云端部署：',True,18),
    ('• 本地 RTX 3060 6GB（不够跑 WavLM 94M + Adapter + DrseCNN）→ 云端 RTX 4090D 24GB',False,15),
    ('• SSH 上传代码 + 数据，conda env 复现 → 完整消融实验约 28 分钟，三 seed 结果自动记录',False,15),
])
add_pn(s,16)

# ═══════════════════════════════════════════
# SLIDE 17 — Full Ablation Results
# ═══════════════════════════════════════════
s=blank_slide(); light_bg(s); section_bar(s,'')
add_text(s,Inches(0.8),Inches(0.3),Inches(11),Inches(0.8),'Phase 3 — 全消融实验结果',fs=36,color=BLACK,bold=True)

add_table(s,Inches(0.8),Inches(1.3),
    [Inches(2.2),Inches(1.2),Inches(1.2),Inches(1.5),Inches(1.5),Inches(1.5)],
    ['实验','Adapter','Pooling','Backbone','Val','Test'],
    [['A1 基线','✗','mean','frozen','85.68%','78.34%'],
     ['A2 统计先验','✓ stat','mean','frozen','86.32%','79.25%'],
     ['A2b 随机初始化','✓ rand','mean','frozen','86.40%','79.51%'],
     ['B3 Adapter+Prosody','✓ rand','prosody','frozen','88.52%','81.24%'],
     ['A3 Full Fine-tune','✓ rand','prosody','unfrozen','88.52%','81.96%']])
add_img(s,charts['module_contrib'],Inches(0.3),Inches(3.8),Inches(6.2))

add_ml(s,Inches(6.8),Inches(3.8),Inches(5.5),Inches(3),[
    ('模块贡献分解（Test）：',True,18),('',False,4),
    ('+ WavLM（SSL 帧级特征）',False,15),
    ('  35% → 78.34% = +43.34pp ← 最大贡献',False,15),('',False,4),
    ('+ Adapter（分布校准）',False,15),
    ('  78.34% → 79.51% = +1.17pp',False,15),('',False,4),
    ('+ Prosody Pooling（韵律引导）',False,15),
    ('  79.51% → 81.24% = +1.73pp',False,15),('',False,4),
    ('+ Full Fine-tune',False,15),
    ('  81.24% → 81.96% = +0.72pp',False,15),
])
add_pn(s,17)

# ═══════════════════════════════════════════
# SLIDE 18 — Three-Phase Comparison
# ═══════════════════════════════════════════
s=blank_slide(); light_bg(s); section_bar(s,'')
add_text(s,Inches(0.8),Inches(0.3),Inches(11),Inches(0.8),'从86%到35%到82%——完整演进',fs=36,color=BLACK,bold=True)

add_img(s,charts['3phase'],Inches(0.3),Inches(1.2),Inches(5.5))
add_table(s,Inches(6.8),Inches(1.2),[Inches(1.8),Inches(1.6),Inches(1.6),Inches(1.9)],
    ['维度','原始论文','修复后','当前方案'],
    [['特征','162-dim mean','同左','WavLM (T×768)'],
     ['时间维度','✗','✗','✓'],
     ['Speaker划分','无(泄露)','70/30','60/20/20'],
     ['增强约束','无','无','分布约束'],
     ['核心创新','网络结构','无','三模块'],
     ['准确率','86.82%','~35%','81.96%'],
     ['可信度','0','—','✓']])

add_text(s,Inches(0.8),Inches(6.0),Inches(11),Inches(0.6),
    '关键信息：86.82% → 35% → 81.96%，数值回到原点但性质完全不同',fs=20,color=RED,bold=True,align=PP_ALIGN.CENTER)
add_pn(s,18)

# ═══════════════════════════════════════════
# SLIDE 19 — Current Research Landscape (NEW)
# ═══════════════════════════════════════════
s=blank_slide(); light_bg(s); section_bar(s,'')
add_text(s,Inches(0.8),Inches(0.3),Inches(11),Inches(0.8),'这个方向在现有研究中处于什么位置？',fs=36,color=BLACK,bold=True)

add_ml(s,Inches(0.8),Inches(1.3),Inches(6),Inches(5.5),[
    ('相关领域研究现状（2024-2025）：',True,20),('',False,4),
    ('1. 成人→儿童SER迁移',False,17),
    ('   Lesyk et al.(2024)用Wav2Vec从成人IEMOCAP迁移到儿童AIBO数据，首次研究年龄/性别对情绪映射的影响。但仅用线性探针，无专门适配结构。',False,13),
    ('',False,4),
    ('2. SER中的PEFT适配器',False,17),
    ('   Li et al.(ICASSP 2024)系统研究LoRA等参数高效微调在跨域SER中的应用，证明适配器可超越全参数微调。但针对的是"域适配"而非"分布偏移"。',False,13),
    ('',False,4),
    ('3. 分布适配在SER中的应用',False,17),
    ('   ICASSP 2025 "Emotional Distribution Alignment"通过欠采样对齐训练/目标域的标签分布，非特征空间的显式校准。',False,13),
    ('   Interspeech 2024 "Hierarchical Distribution Adaptation"做帧/段/句三层分布适配，但非针对儿童语音。',False,13),
],color=BLACK)

add_ml(s,Inches(7.2),Inches(1.3),Inches(5.5),Inches(5.5),[
    ('本方案的研究空白（Research Gap）：',True,20),('',False,6),
    ('核心空白 1',True,18),
    ('WavLM + 儿童SER的组合在已发表文献中未发现。WavLM被用于儿童ASR和成人SER，但两者交叉的应用是空白的。',False,14),('',False,6),
    ('核心空白 2',True,18),
    ('"分布约束增强"概念在SER中未见。现有增强研究关注"如何生成更多数据"，而非"在什么约束空间下生成才是合法的"。',False,14),('',False,6),
    ('核心空白 3',True,18),
    ('"负面实验"（显示分布偏移定量预测分类退化）作为核心叙事手法，在SER论文中罕见——大多数论文只报告"加了某模块提升了X%"。',False,14),('',False,6),
    ('→ 三个空白叠加，构成了论文的差异性创新点。',True,16),
],color=BLACK)
add_pn(s,19)

# ═══════════════════════════════════════════
# SLIDE 20 — Old vs New Paper Direction (NEW)
# ═══════════════════════════════════════════
s=blank_slide(); light_bg(s); section_bar(s,'')
add_text(s,Inches(0.8),Inches(0.3),Inches(11),Inches(0.8),'论文方向对比：从"拼网络"到"挖分布"',fs=36,color=BLACK,bold=True)

add_table(s,Inches(0.5),Inches(1.3),
    [Inches(2.2),Inches(5.3),Inches(5.3)],
    ['维度','原始论文方向（2025）','当前论文方向（2026）'],
    [['核心叙事','"我们提出了一个新的网络结构DrseCNN"','"儿童语音存在系统性分布偏移，我们设计了三模块来应对这个偏移的每个环节"'],
     ['创新来源','网络模块的组合（ResNet+SE+多stage）','儿童语音统计分布的定量分析'],
     ['实验范型','消融实验（证明每个组件有效）','负面实验（无约束增强 → 分布偏移 → 准确率下降）+ 消融'],
     ['特征手段','手工声学特征 + mean pooling','SSL帧级特征 + Prosody引导池化（保留时间维）'],
     ['增强定位','提高准确率的工具','分布偏移的定量验证手段（FD ∝ 1/Accuracy）'],
     ['对审稿人的说服力','弱：ResNet+SE的组合不是真正的创新','中强：每个模块由分布分析驱动 + 负面实验验证 + 研究空白明确'],
     ['目标贡献','"更好的模型"','"从儿童语音分布出发，重新思考SER的每个环节"']])

add_ml(s,Inches(0.8),Inches(5.5),Inches(11.5),Inches(1.7),[
    ('关键转变：',True,18),('',False,4),
    ('• 原始论文的核心是"网络设计"——读者会问"为什么这个设计更好？"（很难自证）',False,15),
    ('• 当前论文的核心是"分布偏移分析"——读者会看到"儿童≠成人 → 三个具体环节受影响 → 每个环节有定量实验支撑"（逻辑链条完整）',False,15),
    ('• 叙事重心从"我做的模型最强"变成"我发现了一个被忽视的问题，并提出了结构化解决方案"——论证难度更低，说服力更高。',False,15),
])
add_pn(s,20)

# ═══════════════════════════════════════════
# SLIDE 21 — Evidence Chain
# ═══════════════════════════════════════════
s=blank_slide(); light_bg(s); section_bar(s,label='第五部分：总结与下一步')
add_text(s,Inches(0.8),Inches(0.3),Inches(11),Inches(0.8),'已产出的论文核心素材',fs=36,color=BLACK,bold=True)

add_ml(s,Inches(0.8),Inches(1.3),Inches(5.8),Inches(5.5),[
    ('五大核心证据：',True,22),('',False,6),
    ('1. 分布偏移存在：Frechet Distance=0.37，scale偏离8%',False,17),
    ('',False,4),
    ('2. FD ∝ 1/Accuracy：分布偏移可定量预测分类退化',False,17),
    ('   → 四组增强对比实验，FD从0到11.99，Acc从86.42%到43.25%',False,13),
    ('',False,4),
    ('3. 统计先验 < 随机初始化：证明Adapter结构的必要性',False,17),
    ('   → 偏移是非线性的、层相关的，不是简单线性平移',False,13),
    ('',False,4),
    ('4. 儿童情感对增强高度脆弱：',False,17),
    ('   约束参数下仍-19.57pp → 说明增强是高风险操作',False,13),
    ('',False,4),
    ('5. 轻量高效：~900K参数捕获~94M backbone的大部分收益',False,17),
],color=BLACK)

add_ml(s,Inches(7),Inches(1.3),Inches(5.5),Inches(5.5),[
    ('已产出图表：',True,22),('',False,6),
    ('📊 FD vs Accuracy 严格负相关曲线',False,16),
    ('📊 全消融实验主表（5实验×3seeds）',False,16),
    ('📊 模块贡献分解图（从35%到81.96%）',False,16),
    ('📊 三阶段核心对比表',False,16),
    ('📊 跨语言迁移对比（语言的"分布鸿沟"）',False,16),('',False,8),
    ('待补（论文 Statistical Analysis 章节）：',True,18),
    ('• UMAP/t-SNE 可视化：儿童vs成人SSL隐空间分布',False,15),
    ('• 儿童 vs 成人声学参数对比图（F0/formant/语速）',False,15),
    ('• 论文三模块架构总图',False,15),
])
add_pn(s,21)

# ═══════════════════════════════════════════
# SLIDE 22 — Paper Structure
# ═══════════════════════════════════════════
s=blank_slide(); light_bg(s); section_bar(s,'')
add_text(s,Inches(0.8),Inches(0.3),Inches(11),Inches(0.8),'论文框架',fs=36,color=BLACK,bold=True)

add_box(s,Inches(1),Inches(1.3),Inches(11),Inches(5.3),fill=LIGHT_GRAY)
add_ml(s,Inches(1.5),Inches(1.5),Inches(10),Inches(5),[
    ('1. Introduction',True,18),('   儿童≠成人，现有SER忽视分布偏移 → 本文提出分布驱动三模块框架',False,15),('',False,6),
    ('2. Statistical Analysis',True,18),('   F0/formant/语速定量对比 + SSL隐空间儿童vs成人分布可视化',False,15),('',False,6),
    ('3. Method',True,18),
    ('   3.1  Distribution Calibration Adapter（特征空间分布校准）',False,15),
    ('   3.2  Temporal Importance Pooling（韵律引导的帧级池化）',False,15),
    ('   3.3  Distribution-Constrained Augmentation（分布约束增强+负面实验）',False,15),
    ('   3.4  Classifier: DrseNet（自然汇聚，非核心创新，一句话带过）',False,15),('',False,6),
    ('4. Experiments',True,18),('   负面实验（分布偏移→分类退化）+ 全消融 + 跨语言迁移（语言=分布偏移的极端维度）',False,15),('',False,6),
    ('5. Related Work & Conclusion',True,18),
])
add_text(s,Inches(1),Inches(6.8),Inches(10),Inches(0.5),'时间规划：2周论文撰写',fs=16,color=BLACK)
add_pn(s,22)

# ═══════════════════════════════════════════
# SLIDE 23 — Remaining Work
# ═══════════════════════════════════════════
s=blank_slide(); light_bg(s); section_bar(s,'')
add_text(s,Inches(0.8),Inches(0.3),Inches(11),Inches(0.8),'待补实验与风险',fs=36,color=BLACK,bold=True)

add_ml(s,Inches(0.8),Inches(1.3),Inches(5.5),Inches(4.5),[
    ('可选补充实验（非必需，但有加分）：',True,20),('',False,6),
    ('• 跨语言迁移（ENGLISH↔TELUGU）✅已初步完成',False,16),
    ('  X1: 19.68%, X2: 28.15% → 确认语言鸿沟',False,13),('',False,6),
    ('• emotion2vec+ / WavLM Large 对比',False,16),
    ('  emotion2vec+ Large→22.45%（效果差，待排查FunASR问题）',False,13),('',False,6),
    ('• UMAP/t-SNE 分布可视化',False,16),
    ('  论文 Statistical Analysis 章节需要',False,13),
    ('• 更多跨语料验证（如非BESD数据集）',False,16),
],color=BLACK)

add_ml(s,Inches(7),Inches(1.3),Inches(5.5),Inches(4.5),[
    ('已解决的问题：',True,20),('',False,6),
    ('✅ emotion2vec 加载(FunASR+ModelScope)',False,16),
    ('✅ 6GB VRAM瓶颈(→云端RTX 4090D 24GB)',False,16),
    ('✅ 数据泄露(两次修复:70/30→60/20/20)',False,16),
    ('✅ 实验日志自动记录、可追溯',False,16),
    ('✅ WavLM本地缓存(无需联网) + 云端快照路径',False,16),('',False,8),
    ('当前优先级：',True,20),
    ('• 论文撰写是第一优先级',False,17),
    ('• 跨语言实验已完成，可作为论文5.5节',False,16),
    ('• 如有审稿人要求，再补emotion2vec+或跨语料实验',False,16),
],color=BLACK)
add_pn(s,23)

# ═══════════════════════════════════════════
# SLIDE 24 — Summary
# ═══════════════════════════════════════════
s=blank_slide(); dark_bg(s)
add_text(s,Inches(0.8),Inches(0.3),Inches(11),Inches(1),'总结',fs=42,color=WHITE,bold=True,align=PP_ALIGN.CENTER)

add_text(s,Inches(0.8),Inches(1.5),Inches(11),Inches(0.6),'项目的两个关键转折点',fs=24,color=PRIMARY,bold=True,align=PP_ALIGN.CENTER)
add_table(s,Inches(1),Inches(2.3),[Inches(1.5),Inches(3.5),Inches(6.5)],
    ['转折','触发','导致的变化'],
    [['第一','面试时CTO质疑增强策略','86.82%→35%，发现真实基线'],
     ['第二','入职时CTO纠正我的SSL偏见','35%→81.96%，形成三模块框架']])

add_text(s,Inches(0.8),Inches(4.3),Inches(11),Inches(0.6),'三句话概括',fs=24,color=PRIMARY,bold=True,align=PP_ALIGN.CENTER)
for i,(idx,txt) in enumerate([
    ('1.','原始方案 86.82% 由数据泄露导致，不可用'),
    ('2.','修复后 162-dim mean-pooled 真实上限 ~35%，确认特征表示是瓶颈'),
    ('3.','当前方案以"儿童分布偏移"为主线，三模块驱动，严格划分下达81.96%'),
]):
    add_box(s,Inches(1.5),Inches(5.0+i*0.65),Inches(10),Inches(0.55),fill=RGBColor(0x2A,0x2A,0x4E))
    add_text(s,Inches(1.8),Inches(5.05+i*0.65),Inches(9.5),Inches(0.45),f'{idx} {txt}',fs=16,color=WHITE)

add_text(s,Inches(0.8),Inches(7.0),Inches(11),Inches(0.5),
    '核心贡献：不是"更强的模型"，而是"从儿童语音分布出发，重新思考SER的每个环节应该怎么做"',
    fs=14,color=GRAY,align=PP_ALIGN.CENTER)
add_pn(s,24)

# ═══════════════════════════════════════════
# SLIDE 25 — Q&A
# ═══════════════════════════════════════════
s=blank_slide(); dark_bg(s)
add_text(s,Inches(2),Inches(2.5),Inches(9),Inches(1.5),'谢谢！欢迎提问',fs=52,color=WHITE,bold=True,align=PP_ALIGN.CENTER)
line=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(4.5),Inches(4.2),Inches(4),Inches(0.03))
line.fill.solid(); line.fill.fore_color.rgb=PRIMARY; line.line.fill.background()
add_text(s,Inches(2),Inches(4.5),Inches(9),Inches(0.6),'GitHub: Xuzc317/child-speech-emotion-recognition',fs=16,color=GRAY,align=PP_ALIGN.CENTER)
add_text(s,Inches(2),Inches(5.1),Inches(9),Inches(0.5),'项目路径: D:/大学/论文/儿童语音情绪识别/新方案-分布驱动儿童SER',fs=14,color=GRAY,align=PP_ALIGN.CENTER)
add_text(s,Inches(2),Inches(5.5),Inches(9),Inches(0.5),'感谢导师指导',fs=14,color=GRAY,align=PP_ALIGN.CENTER)

# ═══════════════════════ Save ═══════════════════════
prs.save(OUTPUT)
print(f'PPT saved: {OUTPUT}')
print(f'Total slides: {len(prs.slides)}')
